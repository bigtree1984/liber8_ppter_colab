"""
converter.py  —  liber8_ppter A4H2P コンバーター
a4h2p_spec.md（liber8_ppter版）準拠

【ローカル実行】
  python converter.py input/slides.html
  → output/liber8_ppter_YYMMDDHHMMSS.pptx

【FastAPI組み込み時】
  from converter import convert_html
  pptx_bytes = convert_html(html_str, template_path)
"""

import re
import sys
import io
from pathlib import Path
from datetime import datetime, timezone, timedelta
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from PIL import Image as PILImage
from lxml import etree

# ════════════════════════════════════════════════════════
#  定数
# ════════════════════════════════════════════════════════

# HTMLキャンバスサイズ
HTML_W = 1920
HTML_H = 1080

# PPTX スライドサイズ（25.4cm × 14.29cm）
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# 座標変換係数: HTML 1920px = PPTX 10inch
COORD_SCALE = 10 / HTML_W  # 0.005208...

# フォントサイズ換算係数（liber8_ppter実測: HTML px × 0.375 = PPTX pt）
FONT_SCALE = 0.375

# POTXテンプレートパス（ローカル実行時のデフォルト）
DEFAULT_TEMPLATE = Path(__file__).parent.parent / "templates" / "libre8_ppter_a.potx"

# 出力ディレクトリ（ローカル実行時）
OUTPUT_DIR = Path(__file__).parent.parent / "output"

# タイムゾーン（JST）
JST = timezone(timedelta(hours=9))

# レイアウト名（POTX定義と一致させること）
LAYOUT_NAMES = ["01_title", "02_structure_dark", "03_structure_light", "04_body"]


# ════════════════════════════════════════════════════════
#  ユーティリティ
# ════════════════════════════════════════════════════════

def parse_style(style_str: str) -> dict:
    """style属性文字列 → dict"""
    result = {}
    for item in (style_str or '').split(';'):
        item = item.strip()
        if ':' in item:
            k, v = item.split(':', 1)
            result[k.strip()] = v.strip()
    return result

def px(val_str: str) -> float:
    """'80px' → 80.0"""
    return float(str(val_str).replace('px', '').strip())

def coord(px_val: float):
    """HTML px → EMU（PPTX座標）"""
    return Inches(px_val * COORD_SCALE)

def hex_to_rgb(val: str) -> RGBColor:
    """'#RRGGBB' → RGBColor"""
    val = val.strip().lstrip('#')
    return RGBColor(int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))

def get_bg_color(style_dict: dict) -> RGBColor | None:
    """background / background-color から RGBColor を取得"""
    val = style_dict.get('background') or style_dict.get('background-color') or ''
    val = val.strip()
    if val.startswith('#'):
        return hex_to_rgb(val)
    return None

def font_pt(css_px_str: str) -> float:
    """'59px' → 22.125pt（× 0.375）"""
    return float(str(css_px_str).replace('px', '').strip()) * FONT_SCALE

def to_align(val: str) -> PP_ALIGN:
    if 'center' in val: return PP_ALIGN.CENTER
    if 'right'  in val: return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT

def is_bold(val: str) -> bool:
    try:    return int(val) >= 700
    except: return val in ('bold', 'bolder')

def resolve_font(css_font: str) -> str | None:
    """CSS font-family → PPTXフォント名。汎用ファミリー名はNoneを返す"""
    if not css_font:
        return None
    first = css_font.split(',')[0].strip().strip("'\"")
    if first.lower() in ('sans-serif', 'serif', 'monospace', 'cursive', 'fantasy'):
        return None
    return first

def clear_shadow(shape) -> None:
    """テーマ由来のシャドウを無効化（空のeffectLstで上書き）"""
    sp_pr = shape._element.spPr
    if sp_pr.find(qn('a:effectLst')) is None:
        etree.SubElement(sp_pr, qn('a:effectLst'))

def get_layout(prs: Presentation, name: str):
    """レイアウト名でレイアウトを取得。見つからなければblank(6)にフォールバック"""
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == name:
            return layout
    print(f'  ⚠ レイアウト"{name}"が見つかりません → blankにフォールバック')
    return prs.slide_masters[0].slide_layouts[6]

def output_filename() -> str:
    """liber8_ppter_YYMMDDHHMMSS.pptx"""
    ts = datetime.now(JST).strftime("%y%m%d%H%M%S")
    return f"liber8_ppter_{ts}.pptx"


# ════════════════════════════════════════════════════════
#  ハンドラー
# ════════════════════════════════════════════════════════

def handle_page_title(slide, el, _base_dir=None) -> None:
    """
    data-pptx-type="page-title"
    04_bodyレイアウトのタイトルエリアにテキストボックスとして配置。
    PHは全削除済みのためtextboxで直接書き込む。
    位置・サイズはPOTX実測値固定（左:0.50cm 上:0.00cm 幅:6.27cm 高:0.94cm）。
    """
    # POTX実測値をEMUで指定（px換算: 0.50cm=38px / 0.00cm=0px / 6.27cm=474px / 0.94cm=71px）
    from pptx.util import Cm
    text = el.get_text(strip=True)
    tb = slide.shapes.add_textbox(Cm(0.50), Cm(0.00), Cm(6.27), Cm(0.94))
    tf = tb.text_frame
    tf.word_wrap     = False
    tf.margin_left   = 0
    tf.margin_right  = 0
    tf.margin_top    = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment    = PP_ALIGN.LEFT
    p.line_spacing = 1.0
    r = p.add_run()
    r.text           = text
    r.font.size      = Pt(22)        # type-title: 22pt
    r.font.bold      = True
    r.font.color.rgb = RGBColor(0x00, 0x88, 0xC8)  # accent1 小林ブルー
    r.font.name      = 'Noto Sans JP'
    tf.auto_size     = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    print(f'    [page-title] → "{text}"')


def handle_text(slide, el, _base_dir=None) -> None:
    """data-pptx-type="text" → テキストボックス"""
    s = parse_style(el.get('style', ''))

    tb = slide.shapes.add_textbox(
        coord(px(s['left'])),  coord(px(s['top'])),
        coord(px(s['width'])), coord(px(s['height'])))
    tf = tb.text_frame
    tf.word_wrap     = False  # 折り返しOFF（auto_sizeと組み合わせ）
    tf.margin_left   = 0
    tf.margin_right  = 0
    tf.margin_top    = 0
    tf.margin_bottom = 0

    color     = hex_to_rgb(s['color']) if s.get('color', '').startswith('#') else RGBColor(0, 0, 0)
    fs_pt     = font_pt(s.get('font-size', '48px'))
    bold      = is_bold(s.get('font-weight', '400'))
    align     = to_align(s.get('text-align', 'left'))
    font_name = resolve_font(s.get('font-family', ''))

    # <br> で行分割
    inner = el.decode_contents()
    lines = re.split(r'<br\s*/?>', inner, flags=re.IGNORECASE)

    for i, line_html in enumerate(lines):
        text = BeautifulSoup(line_html, 'html.parser').get_text(strip=True)
        if not text:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment    = align
        p.line_spacing = 1.0  # シングルスペース固定
        r = p.add_run()
        r.text           = text
        r.font.size      = Pt(fs_pt)
        r.font.bold      = bold
        r.font.color.rgb = color
        if font_name:
            r.font.name = font_name

    # ③ テキスト流し込み後にauto_size設定
    # CENTER/RIGHTはテキストボックスの左端基準でサイズ拡張されるため位置ずれなし
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT


def handle_rect(slide, el, _base_dir=None) -> None:
    """data-pptx-type="rect" / "line" → 矩形シェイプ"""
    s = parse_style(el.get('style', ''))
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        coord(px(s['left'])),  coord(px(s['top'])),
        coord(px(s['width'])), coord(px(s['height'])))
    color = get_bg_color(s)
    if color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    clear_shadow(shp)

    if el.has_attr('data-pptx-shadow'):
        _apply_shadow(shp, el['data-pptx-shadow'])


def handle_line(slide, el, _base_dir=None) -> None:
    """data-pptx-type="line" → rectと同処理"""
    handle_rect(slide, el)


def handle_oval(slide, el, _base_dir=None) -> None:
    """data-pptx-type="oval" → 楕円シェイプ"""
    s = parse_style(el.get('style', ''))
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        coord(px(s['left'])),  coord(px(s['top'])),
        coord(px(s['width'])), coord(px(s['height'])))
    color = get_bg_color(s)
    if color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    clear_shadow(shp)


def handle_triangle(slide, el, _base_dir=None) -> None:
    """
    data-pptx-type="triangle" → 等辺三角形シェイプ
    data-pptx-rotation: 0=▲上 / 90=▶右 / 180=▼下 / 270=◀左
    clip-path（HTML描写用）はConverterが無視する。
    """
    s = parse_style(el.get('style', ''))
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        coord(px(s['left'])),  coord(px(s['top'])),
        coord(px(s['width'])), coord(px(s['height'])))
    color = get_bg_color(s)
    if color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    clear_shadow(shp)
    rotation = float(el.get('data-pptx-rotation', '0'))
    if rotation:
        shp.rotation = rotation


def make_handle_placeholder(images: dict):
    """
    imagesを束縛したhandle_placeholder関数を返す。
    images = {"ラベル名": Path("/tmp/xxxx.jpg")}
    """
    def handle_placeholder(slide, el, _base_dir=None) -> None:
        """
        data-pptx-type="placeholder"
        imagesにlabelが存在すれば画像を埋め込み、なければグレー矩形。
        """
        s     = parse_style(el.get('style', ''))
        label = el.get('data-placeholder-label', '画像')
        hint  = el.get('data-placeholder-hint', '')
        x = coord(px(s['left']))
        y = coord(px(s['top']))
        w = px(s['width'])
        h = px(s['height'])

        img_path = images.get(label)
        if img_path and Path(img_path).exists():
            print(f'    [placeholder] "{label}" → 画像を埋め込み')
            add_picture_cover(slide, img_path, x, y, coord(w), coord(h))
        else:
            print(f'    [placeholder] "{label}" → グレー矩形（ヒント: {hint[:40]}）')
            shp = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, coord(w), coord(h))
            shp.fill.solid()
            shp.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            shp.line.fill.background()

    return handle_placeholder


def handle_image(slide, el, base_dir: Path) -> None:
    """data-pptx-type="image" → 画像埋め込み（URLはダウンロード）"""
    s   = parse_style(el.get('style', ''))
    # src is on the div itself or on a child <img> tag
    img_tag = el.find('img') if hasattr(el, 'find') else None
    src = el.get('src', '') or (img_tag.get('src', '') if img_tag else '')
    x   = coord(px(s['left']))
    y   = coord(px(s['top']))
    w   = px(s['width'])
    h   = px(s['height'])

    # URL画像はダウンロード
    if src.startswith('http'):
        import urllib.request, tempfile
        try:
            suffix = Path(src.split('?')[0]).suffix or '.jpg'
            tmp = Path(tempfile.mktemp(suffix=suffix))
            urllib.request.urlretrieve(src, tmp)
            img_path = tmp
        except Exception as e:
            print(f'    ⚠ URL画像ダウンロード失敗: {e}')
            img_path = None
    else:
        img_path = (base_dir / src).resolve()

    if img_path is None or not img_path.exists():
        print(f'    ⚠ 画像が見つかりません: {src} → グレー矩形で代替')
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, coord(w), coord(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        shp.line.fill.background()
        return

    add_picture_cover(slide, img_path, x, y, coord(w), coord(h))


# ════════════════════════════════════════════════════════
#  画像配置（カバーフィット・クロップ）
# ════════════════════════════════════════════════════════

def add_picture_cover(slide, img_path: Path, x, y, w_emu: int, h_emu: int) -> None:
    """
    画像をcover fitで配置する（縦横比を保ちながら枠を目一杯埋め、はみ出た部分をクロップ）。
    """
    # 元画像のサイズを取得
    with PILImage.open(str(img_path)) as img:
        img_w, img_h = img.size

    # cover fit のスケール計算（長辺基準）
    scale_w = w_emu / img_w
    scale_h = h_emu / img_h
    scale   = max(scale_w, scale_h)

    scaled_w = int(img_w * scale)
    scaled_h = int(img_h * scale)

    # クロップ割合を計算（単位: 100000 = 100%、中央揃え）
    crop_x = (scaled_w - w_emu) / 2
    crop_y = (scaled_h - h_emu) / 2

    l = int(crop_x / scaled_w * 100000)
    t = int(crop_y / scaled_h * 100000)
    r = int((scaled_w - w_emu - crop_x) / scaled_w * 100000)
    b = int((scaled_h - h_emu - crop_y) / scaled_h * 100000)

    # 枠サイズで配置（python-pptxはここで画像をストレッチする）
    pic = slide.shapes.add_picture(str(img_path), x, y, w_emu, h_emu)
    pic_elem = pic._element

    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    # p:blipFill を取得（p:名前空間）
    blip_fill = pic_elem.find(f'{{{p_ns}}}blipFill')
    if blip_fill is None:
        blip_fill = pic_elem.find(f'.//{{{a_ns}}}blipFill')
    if blip_fill is None:
        return

    # blipを取得
    blip = blip_fill.find(f'{{{a_ns}}}blip')

    # blipFillを再構築: blip → srcRect → stretch
    blip_fill[:] = []
    if blip is not None:
        blip_fill.append(blip)

    src_rect = etree.SubElement(blip_fill, f'{{{a_ns}}}srcRect')
    src_rect.set('l', str(max(0, l)))
    src_rect.set('t', str(max(0, t)))
    src_rect.set('r', str(max(0, r)))
    src_rect.set('b', str(max(0, b)))

    stretch = etree.SubElement(blip_fill, f'{{{a_ns}}}stretch')
    etree.SubElement(stretch, f'{{{a_ns}}}fillRect')


# ════════════════════════════════════════════════════════
#  シャドウ補助
# ════════════════════════════════════════════════════════

def _apply_shadow(shape, shadow_attr: str) -> None:
    """
    data-pptx-shadow="preset1" → PPTプリセット1シャドウを適用
    #000000・透明度60%・ぼかし4pt・角度45°・距離3pt
    """
    sp_pr   = shape._element.spPr
    eff_lst = sp_pr.find(qn('a:effectLst'))
    if eff_lst is None:
        eff_lst = etree.SubElement(sp_pr, qn('a:effectLst'))
    for old in eff_lst.findall(qn('a:outerShdw')):
        eff_lst.remove(old)

    outer = etree.SubElement(eff_lst, qn('a:outerShdw'))
    outer.set('blurRad',       str(int(Pt(4))))     # 4pt ぼかし
    outer.set('dist',          str(int(Pt(3))))     # 3pt 距離
    outer.set('dir',           str(45 * 60000))     # 45°
    outer.set('algn',          'tl')
    outer.set('rotWithShape',  '0')

    srgb = etree.SubElement(outer, qn('a:srgbClr'))
    srgb.set('val', '000000')
    alpha_el = etree.SubElement(srgb, qn('a:alpha'))
    alpha_el.set('val', str(int(0.4 * 100000)))     # 透明度60% = alpha40%


# ════════════════════════════════════════════════════════
#  ハンドラーマップ
# ════════════════════════════════════════════════════════

HANDLERS = {
    'page-title':  handle_page_title,
    'text':        handle_text,
    'rect':        handle_rect,
    'line':        handle_line,
    'oval':        handle_oval,
    'triangle':    handle_triangle,
    'placeholder': None,  # convert_html()内でmake_handle_placeholder()で動的生成
    'image':       handle_image,
    'ignore':      lambda *_: None,
}


# ════════════════════════════════════════════════════════
#  コア変換関数（FastAPI・CLIで共用）
# ════════════════════════════════════════════════════════

def convert_html(html_str: str, template_path: Path, base_dir: Path = Path('.'), images: dict = None) -> bytes:
    """
    A4H2P形式のHTML文字列をPPTXに変換してbytesで返す。
    FastAPIからはこの関数を呼び出す。

    Args:
        html_str:      A4H2P形式のHTML文字列
        template_path: POTXテンプレートのパス
        base_dir:      画像ファイルの基点ディレクトリ（inputフォルダ）
        images:        placeholderに埋め込む画像 {"ラベル名": Path}

    Returns:
        PPTXファイルのbytes
    """
    if images is None:
        images = {}
    soup = BeautifulSoup(html_str, 'html.parser')

    # POTXをPPTXとして読み込む（Content_Types書き換え）
    import zipfile, shutil, tempfile
    tmp_pptx = Path(tempfile.mktemp(suffix='.pptx'))
    shutil.copy(template_path, tmp_pptx)
    with zipfile.ZipFile(tmp_pptx, 'r') as zin:
        contents = {name: zin.read(name) for name in zin.namelist()}
    ct = contents['[Content_Types].xml'].decode('utf-8')
    ct = ct.replace(
        'application/vnd.openxmlformats-officedocument.presentationml.template.main+xml',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml'
    )
    contents['[Content_Types].xml'] = ct.encode('utf-8')
    with zipfile.ZipFile(tmp_pptx, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, data in contents.items():
            zout.writestr(name, data)

    prs = Presentation(str(tmp_pptx))
    tmp_pptx.unlink()

    # POTXに含まれるサンプルスライドを全削除
    r_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(f'{{{r_ns}}}id')
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    sections = soup.select('section[data-pptx-layout]')
    print(f'スライド数: {len(sections)}')

    # imagesを束縛したplaceholderハンドラーを生成
    handlers = dict(HANDLERS)
    handlers['placeholder'] = make_handle_placeholder(images)

    for section in sections:
        layout_name = section.get('data-pptx-layout', '04_body')
        label       = section.get('data-label', '(無名)')
        print(f'  → [{layout_name}] {label}')

        layout = get_layout(prs, layout_name)
        slide  = prs.slides.add_slide(layout)

        # ② 全PHを削除（page-titleハンドラーで必要なidx=0は後で書き込む）
        # 02/03: idx=0（セクションタイトル）
        # 04:    idx=0（スライドタイトル）+ idx=10（コンテンツ）
        # いずれもConverterでは使わないため全削除
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        for el in section.select('[data-pptx-type]'):
            ptype   = el.get('data-pptx-type')
            handler = handlers.get(ptype)
            if handler:
                try:
                    handler(slide, el, base_dir)
                except Exception as e:
                    print(f'    ⚠ エラー [{ptype}]: {e}')
            else:
                print(f'    ⚠ 未知のtype="{ptype}" — スキップ')

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    # ダミーフォント経由の2回置換
    # PowerPointが「フォント変更あり」と認識してauto_sizeを発火させる
    DUMMY = b'DUMMY_FONT_LIBR8'
    pptx_bytes = pptx_bytes.replace(b'Noto Sans JP', DUMMY)
    pptx_bytes = pptx_bytes.replace(DUMMY, b'Noto Sans JP')
    print('  フォント2回置換完了')

    return pptx_bytes


# ════════════════════════════════════════════════════════
#  CLI エントリーポイント
# ════════════════════════════════════════════════════════

def main():
    if len(sys.argv) < 2:
        print("使い方: python converter.py input/slides.html [template_path]")
        sys.exit(1)

    html_path     = Path(sys.argv[1])
    template_path = Path(sys.argv[2]) if len(sys.argv) >= 3 else DEFAULT_TEMPLATE

    if not html_path.exists():
        print(f"エラー: HTMLファイルが見つかりません: {html_path}")
        sys.exit(1)
    if not template_path.exists():
        print(f"エラー: テンプレートが見つかりません: {template_path}")
        sys.exit(1)

    html_str = html_path.read_text(encoding='utf-8')
    base_dir = html_path.parent

    print(f"入力  : {html_path}")
    print(f"テンプレ: {template_path}")

    pptx_bytes = convert_html(html_str, template_path, base_dir)

    OUTPUT_DIR.mkdir(exist_ok=True)
    out_path = OUTPUT_DIR / output_filename()
    out_path.write_bytes(pptx_bytes)

    print(f"\n✅ 保存完了: {out_path}")


if __name__ == '__main__':
    main()